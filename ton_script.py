from pptx import Presentation
from pptx.util import Inches
import re

# Lire le fichier Markdown localement
with open("eca1.md", "r", encoding="utf-8") as f:
    content = f.read()

# Séparer le contenu en pages (repère : "--- Page X ---")
pages = re.split(r"## --- Page \d+ ---", content)

# Créer la présentation PowerPoint
prs = Presentation()
layout = prs.slide_layouts[1]  # Titre + contenu

# Fonction pour ajouter une diapositive
def add_slide(title, body):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title.strip() or "Slide"
    slide.placeholders[1].text = body.strip()

# Génération automatique des slides
for page in pages:
    lines = page.strip().split("\n")
    if not lines or all(not line.strip() for line in lines):
        continue
    title = lines[0].strip() if lines[0].strip() else "Contenu"
    body = "\n".join(lines[1:]) if len(lines) > 1 else "Contenu à compléter"
    add_slide(title, body)

# Enregistrement du PowerPoint
prs.save("presentation_auto_eca1.pptx")
print("✅ Présentation générée : presentation_auto_eca1.pptx")
